"""PPTX 处理器"""
from pathlib import Path
from cleaner.state import upsert_metadata, log

BASE_DIR = Path.home() / "Documents" / "knowledge-system"


def process_pptx(file_hash: str, vault_path: str, conn) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        log(conn, file_hash, "processor", "error", "python-pptx 未安装")
        return None

    prs = Presentation(vault_path)
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
            if shape.has_table:
                table = shape.table
                lines.append("| " + " | ".join("—" * 3 for _ in table.columns) + " |")
                for row in table.rows:
                    cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")

        # 备注
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"[备注] {notes}")

    text = "\n".join(lines)
    text_path = BASE_DIR / "processed" / "text" / f"{file_hash}.txt"
    text_path.write_text(text, encoding="utf-8")

    props = prs.core_properties
    meta = {
        "title": props.title or Path(vault_path).stem,
        "author": props.author,
        "slide_count": len(prs.slides),
    }
    upsert_metadata(conn, file_hash, **meta)

    log(conn, file_hash, "processor", "done", f"pptx→text: {len(text)} chars, {len(prs.slides)} slides")
    return {"text_path": str(text_path), "image_paths": [], "meta": meta}
